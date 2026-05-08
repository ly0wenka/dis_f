from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import os

# Створення тимчасової папки для формул
TEMP_DIR = "temp_formulas"
if not os.path.exists(TEMP_DIR):
    os.makedirs(TEMP_DIR)

def formula_to_image(latex_formula, filename, fontsize=14, dpi=120):
    """Перетворює просту LaTeX-формулу на зображення"""
    plt.rcParams['text.usetex'] = False
    plt.rcParams['mathtext.fontset'] = 'cm'
    
    fig, ax = plt.subplots(figsize=(8, 0.7))
    ax.axis('off')
    
    # Додаємо формулу
    ax.text(0.5, 0.5, f"${latex_formula}$", 
            fontsize=fontsize, ha='center', va='center',
            transform=ax.transAxes)
    
    fig.tight_layout(pad=0.5)
    fig.savefig(os.path.join(TEMP_DIR, filename), dpi=dpi, 
                bbox_inches='tight', facecolor='white', edgecolor='none')
    plt.close(fig)
    return os.path.join(TEMP_DIR, filename)

def add_formula_slide(prs, title, formulas, descriptions=None):
    """Додає слайд з формулами"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    left = Inches(0.5)
    top = Inches(1.2)
    
    if descriptions:
        textbox = slide.shapes.add_textbox(left, top, Inches(12), Inches(0.6))
        tf = textbox.text_frame
        tf.text = descriptions
        tf.paragraphs[0].font.size = Pt(11)
    
    y_offset = 0.7 if descriptions else 0.3
    for i, formula in enumerate(formulas):
        img_path = formula_to_image(formula, f"formula_{i}.png", fontsize=13)
        pic = slide.shapes.add_picture(img_path, 
                                       left + Inches(0.2), 
                                       top + Inches(y_offset + i * 0.6),
                                       width=Inches(8), 
                                       height=Inches(0.5))
        # Номер формули
        num_box = slide.shapes.add_textbox(left + Inches(8.5), top + Inches(y_offset + i * 0.6), 
                                           Inches(0.8), Inches(0.5))
        num_frame = num_box.text_frame
        num_frame.text = f"({i+1})"
        num_frame.paragraphs[0].font.size = Pt(10)
        num_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle

def add_content_slide(prs, title, content_lines):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = content_lines[0] if content_lines else ""
    for line in content_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.font.size = Pt(13)

def add_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.3)
    top = Inches(1.2)
    width = Inches(12.5)
    height = Inches(5.5)
    table = slide.shapes.add_table(len(rows)+1, len(headers), left, top, width, height).table
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
    for row_idx, row in enumerate(rows):
        for col_idx, val in enumerate(row):
            cell = table.cell(row_idx+1, col_idx)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(10)

def add_slide_with_image_placeholder(prs, title, text_lines, figure_ref):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = text_lines[0] if text_lines else ""
    for line in text_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.bullet = True
    left = Inches(8)
    top = Inches(2)
    width = Inches(4.5)
    height = Inches(4.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.add_paragraph()
    p.text = f"▶ ТУТ МАЄ БУТИ РИСУНОК:\n{figure_ref}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(128, 128, 128)

# Створення презентації
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ==================== СЛАЙДИ (25 шт.) ====================

# 1. Титульний
add_title_slide(prs, "МОДЕЛІ, МЕТОДИ ТА ІНФОРМАЦІЙНА ТЕХНОЛОГІЯ\nДИСТАНЦІЙНОЇ ІДЕНТИФІКАЦІЇ ПАРАМЕТРІВ ДИНАМІЧНИХ ОБ'ЄКТІВ",
                "Здобувач: Кондратов Олексій Михайлович\nНауковий керівник: д.т.н., проф. Нікуліна Олена Миколаївна\nНТУ «ХПІ» – Харків, 2026")

# 2. Актуальність
add_content_slide(prs, "Актуальність теми", [
    "Підвищення точності, надійності та адаптивності дистанційної ідентифікації параметрів динамічних об'єктів (ДІПДО)",
    "Інтеграція трансформерних детекторів (DETR), методів оптичного потоку (Farneback, Lucas-Kanade, Horn-Schunck, FlowNet, GeoNet)",
    "Використання ансамблевих методів (Bagging, Boosting) для підвищення стійкості",
    "Існуючі системи працюють в обмежених умовах (статична камера, фіксоване освітлення) або потребують значних ресурсів"
])

# 3. Завдання дослідження
add_content_slide(prs, "Завдання дослідження", [
    "1. Аналіз сучасних методів ДІПДО, визначення обмежень",
    "2. Класифікація параметрів динамічних об'єктів та методів їх визначення",
    "3. Розробка моделей інтеграції DETR, оптичного потоку та нейромережевої оцінки глибини",
    "4. Обґрунтування використання Bagging та Boosting",
    "5. Програмна реалізація інформаційної технології",
    "6. Експериментальне тестування 35 комбінацій методів"
])

# 4. СЛАЙД З ФОРМУЛАМИ: Рівняння оптичного потоку
add_formula_slide(prs, "Рівняння оптичного потоку",
                  ["I_x u + I_y v + I_t = 0",
                   "\\frac{\\partial I}{\\partial x} u + \\frac{\\partial I}{\\partial y} v + \\frac{\\partial I}{\\partial t} = 0"],
                  "Основне рівняння оптичного потоку")

# 5. СЛАЙД З ФОРМУЛАМИ: Метод Lucas-Kanade
add_formula_slide(prs, "Метод Lucas-Kanade",
                  ["(A^T A) \\begin{bmatrix} u \\\\ v \\end{bmatrix} = A^T b",
                   "A^T A = \\begin{bmatrix} \\sum I_x^2 & \\sum I_x I_y \\\\ \\sum I_x I_y & \\sum I_y^2 \\end{bmatrix}",
                   "b = - \\begin{bmatrix} \\sum I_x I_t \\\\ \\sum I_y I_t \\end{bmatrix}"],
                  "Система нормальних рівнянь для локального вікна")

# 6. СЛАЙД З ФОРМУЛАМИ: Метод Horn-Schunck
add_formula_slide(prs, "Метод Horn-Schunck",
                  ["E = \\iint \\left[ (I_x u + I_y v + I_t)^2 + \\alpha^2 (\\|\\nabla u\\|^2 + \\|\\nabla v\\|^2) \\right] dx dy",
                   "u^{k+1} = \\bar{u}^k - \\frac{I_x(I_x \\bar{u}^k + I_y \\bar{v}^k + I_t)}{\\alpha^2 + I_x^2 + I_y^2}",
                   "v^{k+1} = \\bar{v}^k - \\frac{I_y(I_x \\bar{u}^k + I_y \\bar{v}^k + I_t)}{\\alpha^2 + I_x^2 + I_y^2}"],
                  "Функціонал енергії та ітераційна схема")

# 7. СЛАЙД З ФОРМУЛАМИ: Метод Farneback
add_formula_slide(prs, "Метод Farneback",
                  ["f(x) \\approx x^T A x + b^T x + c",
                   "f_2(x) = f_1(x - d)",
                   "d = -\\frac{1}{2} A_1^{-1} (b_2 - b_1)"],
                  "Поліноміальне представлення та вектор зміщення")

# 8. СЛАЙД З ФОРМУЛАМИ: Механізм уваги
add_formula_slide(prs, "Механізм уваги (Attention)",
                  ["\\text{Attention}(Q,K,V) = \\text{softmax}\\left(\\frac{QK^T}{\\sqrt{d_k}}\\right) V",
                   "\\text{MultiHead}(Q,K,V) = \\text{Concat}(\\text{head}_1,\\ldots,\\text{head}_h) W^O"],
                  "Багатоголова увага в трансформерах")

# 9. СЛАЙД З ФОРМУЛАМИ: DETR
add_formula_slide(prs, "Функція втрат DETR",
                  ["\\mathcal{L}_{Hungarian} = \\sum_{i=1}^N \\left[ -\\log \\hat{p}(c_i) + \\mathbf{1}_{c_i \\neq \\emptyset} \\mathcal{L}_{box} \\right]",
                   "\\mathcal{L}_{box}(b, \\hat{b}) = \\lambda_{IoU} \\mathcal{L}_{IoU}(b, \\hat{b}) + \\lambda_{L1} \\|b - \\hat{b}\\|_1"],
                  "Угорська втрата (Hungarian loss)")

# 10. СЛАЙД З ФОРМУЛАМИ: Ансамблеві методи
add_formula_slide(prs, "Ансамблеві методи",
                  ["\\hat{y}_{bag}(x) = \\frac{1}{M} \\sum_{i=1}^{M} f_i(x)",
                   "\\hat{y}_{boost}(x) = \\sum_{i=1}^{M} \\alpha_i f_i(x), \\quad \\alpha_i = \\frac{1}{2} \\ln\\left(\\frac{1-\\epsilon_i}{\\epsilon_i}\\right)",
                   "D_{final}(x) = \\begin{cases} D_1(x), & |\\nabla D_1(x)| \\leq \\tau \\\\ D_2(x), & |\\nabla D_1(x)| > \\tau \\end{cases}"],
                  "Bagging (усереднення), Boosting (зважена сума)")

# 11. СЛАЙД З ФОРМУЛАМИ: Кінематичні параметри
add_formula_slide(prs, "Кінематичні параметри об'єкта",
                  ["v = \\frac{1}{|\\Omega|} \\iint_\\Omega \\sqrt{u^2 + v^2} \\, dxdy",
                   "\\theta = \\arctan\\left( \\frac{\\sum \\sin\\theta_i}{\\sum \\cos\\theta_i} \\right)",
                   "x_c = \\frac{x_{\\min} + x_{\\max}}{2}, \\quad y_c = \\frac{y_{\\min} + y_{\\max}}{2}, \\quad z = \\text{median}\\{ D(x,y) \\}"],
                  "Швидкість, напрямок руху, центр та глибина")

# 12. СЛАЙД З ФОРМУЛАМИ: GeoNet
add_formula_slide(prs, "Модель GeoNet",
                  ["I_2(p) = I_1(p + f(p))",
                   "f(p) = T(p) \\cdot \\frac{Z(p) - Z_0}{Z(p)} \\cdot \\text{proj}(p)",
                   "f_{t\\rightarrow s}^{rig}(p_t) = K T_{t\\rightarrow s} D_t(p_t) K^{-1} p_t - p_t"],
                  "Фотометрична узгодженість та геометрична модель")

# 13. СЛАЙД З ФОРМУЛАМИ: Метрики якості
add_formula_slide(prs, "Метрики оцінювання якості",
                  ["\\text{EPE} = \\frac{1}{N} \\sum \\sqrt{(u_i - \\hat{u}_i)^2 + (v_i - \\hat{v}_i)^2}",
                   "\\text{RMSE} = \\sqrt{\\frac{1}{N} \\sum (d_i - \\hat{d}_i)^2}",
                   "\\text{IoU} = \\frac{|A \\cap B|}{|A \\cup B|}, \\quad \\text{DICE} = \\frac{2|A \\cap B|}{|A| + |B|}"],
                  "End-Point Error, RMSE, Intersection over Union, Dice")

# 14. Класифікація параметрів
add_table_slide(prs, "Класифікація параметрів",
                ["Категорія", "Параметри"],
                [["Просторові", "Координати (x,y,z), глибина, розміри, орієнтація"],
                 ["Кінематичні", "Швидкість (v), прискорення (a), напрям (θ), траєкторія"],
                 ["Структурні", "Контур, площа, клас, ознаки форма"]])

# 15. Наукова новизна
add_content_slide(prs, "Наукова новизна", [
    "1. Вперше: метод ДІПДО на основі DETR + оптичного потоку",
    "2. Вперше: структура ІТ з інтеграцією FlowNet, GeoNet, трансформерів",
    "3. Удосконалено метод оцінювання глибини (MiDaS, DPT, GeoNet з Bagging/Boosting)",
    "4. Удосконалено метод ансамблевої агрегації",
    "5. Удосконалено систему метрик для 35 комбінацій",
    "6. Подальший розвиток методу синтезу просторово-часових ознак"
])

# 16. Практичне значення
add_content_slide(prs, "Практичне значення", [
    "Безпека та оборона (відеоспостереження, БПЛА)",
    "Автономний транспорт та логістика (уникнення зіткнень)",
    "Цифровізація та економіка (мобільні пристрої)",
    "Медицина (неінвазивна діагностика)"
])

# 17. Результати DETR
add_table_slide(prs, "Виявлення об'єктів DETR",
                ["ID", "Клас", "Довіра", "x,y", "z", "v (px/кадр)", "θ (град)"],
                [["0", "1", "1,00", "658,5;281,0", "93,0", "6,526", "175,5"],
                 ["1", "3", "0,99", "750,0;70,0", "22,0", "0,017", "-152,4"],
                 ["2", "8", "0,79", "689,0;70,5", "21,0", "0,017", "-143,5"],
                 ["3", "8", "0,74", "688,5;70,5", "21,0", "0,017", "-144,8"],
                 ["4", "1", "1,00", "277,5;264,0", "94,0", "4,592", "-6,5"],
                 ["5", "1", "1,00", "513,5;195,5", "56,0", "1,840", "147,5"]])

# 18. Порівняння оптичного потоку
add_table_slide(prs, "Порівняння методів оптичного потоку",
                ["Метод", "EPE", "AAE", "v об'єкта 0"],
                [["Farneback (ref)", "0,000", "0,000", "6,526"],
                 ["Lucas-Kanade", "0,273", "0,152", "0,013"],
                 ["Horn-Schunck", "0,266", "0,145", "0,325"],
                 ["FlowNet (проксі)", "0,000", "0,000", "6,526"],
                 ["GeoNet (проксі)", "0,000", "0,000", "6,526"]])

# 19. Порівняння глибини
add_table_slide(prs, "Порівняння методів глибини",
                ["Метод", "RMSE", "MAE", "PSNR", "IoU", "Dice"],
                [["MiDaS (ref)", "0,000", "0,000", "∞", "1,000", "1,000"],
                 ["DPT_Large", "12,593", "9,057", "26,128", "0,954", "0,976"],
                 ["GeoNet", "1,660", "0,977", "43,727", "0,995", "0,997"],
                 ["Bagging", "0,713", "0,474", "51,065", "0,996", "0,998"],
                 ["Boosting", "5,088", "0,925", "34,000", "0,981", "0,990"],
                 ["Blur", "106,0", "95,36", "7,623", "0,117", "0,209"],
                 ["Gradient", "124,6", "99,96", "6,221", "0,148", "0,257"]])

# 20. Аналіз результатів
add_content_slide(prs, "Аналіз результатів", [
    "Найкраща точність глибини: Bagging (RMSE=0,713, IoU=0,996)",
    "Щільні методи оптичного потоку дають узгоджені оцінки",
    "Lucas-Kanade придатний лише для локального супроводу",
    "Рекомендована комбінація: DETR + Farneback + Bagging"
])

# 21. Висновки
add_content_slide(prs, "Висновки", [
    "1. Розроблено моделі, методи та ІТ ДІПДО",
    "2. Запропоновано метод DETR+оптичний потік",
    "3. Удосконалено оцінку глибини з Bagging/Boosting",
    "4. Підтверджено перевагу Bagging (RMSE=0,713, IoU=0,996)",
    "5. Система стійка до шумів, оклюзій та зміни освітлення"
])

# 22. Впровадження
add_content_slide(prs, "Впровадження результатів", [
    "НТУ «ХПІ»: «Інтелектуальні системи», «Інтелектуальний аналіз даних»",
    "Метінвест Політехніка: кібербезпека, GUI, захист даних",
    "Радіоастрономічний інститут НАНУ: аналіз сонячних радіосплесків"
])

# 23. Публікації
add_content_slide(prs, "Публікації", [
    "1 стаття у Scopus (IEEE KhPIWeek 2025)",
    "7 статей у фахових виданнях України (категорія Б)",
    "7 матеріалів апробаційного характеру",
    "Всього: 15 наукових праць"
])

# 24. Апробація
add_content_slide(prs, "Апробація роботи", [
    "MicroCAD (2023, 2024, 2025) – Харків",
    "KhPI Week on Advanced Technology (2025)",
    "Конференції молодих вчених НТУ «ХПІ» (2024, 2025)"
])

# 25. Заключний
add_title_slide(prs, "ДЯКУЮ ЗА УВАГУ!", "Кондратов Олексій Михайлович")

# Збереження
prs.save("dissertation_Kondratov_final.pptx")

print("=" * 60)
print("✅ Презентацію створено: dissertation_Kondratov_final.pptx")
print(f"📊 Загальна кількість слайдів: {len(prs.slides)}")
print(f"📐 Тимчасові файли формул у папці: {TEMP_DIR}")
print("=" * 60)

# Очищення тимчасових файлів (за потреби)
# import shutil
# shutil.rmtree(TEMP_DIR)